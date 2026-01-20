#!/usr/bin/env python3
import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt

def parse_markdown(file_path):
    slides = []
    current_slide = None
    
    with open(file_path, 'r', encoding='utf-8') as f:
        for line in f:
            line = line.strip()
            if not line:
                continue
                
            if line.startswith('# '):
                if current_slide:
                    slides.append(current_slide)
                current_slide = {'title': line[2:], 'content': []}
            elif line.startswith('## '):
                if current_slide:
                    current_slide['content'].append({'type': 'subtitle', 'text': line[3:]})
            elif line.startswith('- '):
                if current_slide:
                    current_slide['content'].append({'type': 'bullet', 'text': line[2:]})
            else:
                if current_slide:
                    current_slide['content'].append({'type': 'text', 'text': line})
                    
    if current_slide:
        slides.append(current_slide)
        
    return slides

def create_presentation(slides, output_file):
    prs = Presentation()
    
    for slide_data in slides:
        # Use bullet layout for most slides
        slide_layout = prs.slide_layouts[1] 
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title = slide.shapes.title
        title.text = slide_data['title']
        
        # Add content
        body_shape = slide.shapes.placeholders[1]
        tf = body_shape.text_frame
        
        for item in slide_data['content']:
            if item['type'] == 'bullet':
                p = tf.add_paragraph()
                p.text = item['text']
                p.level = 0
            elif item['type'] == 'text':
                p = tf.add_paragraph()
                p.text = item['text']
                p.level = 0
            elif item['type'] == 'subtitle':
                p = tf.add_paragraph()
                p.text = item['text']
                p.font.bold = True
                
    prs.save(output_file)
    print(f"Presentation saved to {output_file}")

if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: python generate_pptx.py <input.md> <output.pptx>")
        sys.exit(1)
        
    input_file = sys.argv[1]
    output_file = sys.argv[2]
    
    if not os.path.exists(input_file):
        print(f"Error: Input file {input_file} not found")
        sys.exit(1)
        
    try:
        slides = parse_markdown(input_file)
        create_presentation(slides, output_file)
    except ImportError:
        print("Error: python-pptx not installed. Run: pip install python-pptx")
    except Exception as e:
        print(f"Error: {str(e)}")
